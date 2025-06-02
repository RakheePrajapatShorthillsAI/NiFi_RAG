import os
import asyncio
from pptx import Presentation
from docx import Document
from typing import List, Dict, Tuple
from dotenv import load_dotenv
import weaviate
from weaviate.auth import AuthApiKey
from openai import AsyncAzureOpenAI
from httpx_aiohttp import AiohttpTransport
from aiohttp import ClientSession
import openai
import numpy as np
import logging

# Load environment variables
load_dotenv()

# Setup basic logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configure working directory
INPUT_DIR = "./test_data/PreSales"
if not os.path.exists(INPUT_DIR):
    os.makedirs(INPUT_DIR)

# Weaviate Configuration
WEAVIATE_URL = os.getenv("WEAVIATE_URL")
WEAVIATE_API_KEY = os.getenv("WEAVIATE_API_KEY")

# Azure OpenAI Configuration
AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_VERSION = os.getenv("AZURE_OPENAI_API_VERSION", "2024-02-15-preview")
EMBEDDING_MODEL = "text-embedding-ada-002"

def verify_env_variables():
    """Verify required environment variables are set."""
    required_vars = {
        'AZURE_OPENAI_API_KEY': 'Azure OpenAI API key',
        'AZURE_OPENAI_ENDPOINT': 'Azure OpenAI endpoint'
    }
    
    missing_vars = []
    for var, description in required_vars.items():
        if not os.getenv(var):
            missing_vars.append(f"{var} ({description})")
    
    if missing_vars:
        raise ValueError(f"Missing required environment variables: {', '.join(missing_vars)}")

async def generate_embedding(text: str) -> np.ndarray:
    """Generate embedding using Azure OpenAI"""
    try:
        async with AiohttpTransport(client=ClientSession()) as aiohttp_transport:
            httpx_client = openai.DefaultAsyncHttpxClient(transport=aiohttp_transport)
            client = AsyncAzureOpenAI(
                azure_endpoint=AZURE_OPENAI_ENDPOINT,
                api_key=AZURE_OPENAI_API_KEY,
                api_version=AZURE_OPENAI_API_VERSION,
                http_client=httpx_client
            )
            
            logger.info(f"Generating embedding for text: {text[:100]}...")
            response = await client.embeddings.create(
                model=EMBEDDING_MODEL,
                input=[text],
                encoding_format="float"
            )
            
            if not response or not hasattr(response, "data") or not response.data:
                logger.error("Empty or invalid response from OpenAI embedding API")
                return None
                
            return np.array(response.data[0].embedding)
            
    except Exception as e:
        logger.error(f"Error generating embedding: {str(e)}")
        return None

def extract_ppt_content(ppt_file: str) -> List[str]:
    """Extract text content from each slide."""
    prs = Presentation(ppt_file)
    slides_content = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        
        # Add slide number as header
        slide_text.append(f"[Slide {slide_num}]")
        
        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    # Add indentation to distinguish content
                    slide_text.extend("  " + line for line in text.split("\n") if line.strip())
        
        if len(slide_text) > 1:  # Only add if there's content beyond the slide number
            content = "\n".join(slide_text)
            slides_content.append(content)
    
    return slides_content

def extract_docx_content(docx_file: str) -> List[str]:
    """Extract text content from DOCX file."""
    doc = Document(docx_file)
    paragraphs_content = []
    current_chunk = []
    current_size = 0
    max_chunk_size = 2000  # Maximum characters per chunk
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
            
        # If adding this paragraph would exceed chunk size, save current chunk
        if current_size + len(text) > max_chunk_size and current_chunk:
            paragraphs_content.append("\n".join(current_chunk))
            current_chunk = []
            current_size = 0
            
        current_chunk.append(text)
        current_size += len(text)
    
    # Add any remaining content
    if current_chunk:
        paragraphs_content.append("\n".join(current_chunk))
    
    return paragraphs_content

def get_input_files() -> List[Tuple[str, str]]:
    """Get all PPT and DOCX files from input directory."""
    input_files = []
    for filename in os.listdir(INPUT_DIR):
        if filename.endswith(('.ppt', '.pptx', '.docx')):
            file_path = os.path.join(INPUT_DIR, filename)
            file_type = 'ppt' if filename.endswith(('.ppt', '.pptx')) else 'docx'
            input_files.append((file_path, file_type))
    return input_files

async def process_and_upload_file(file_path: str, file_type: str, weaviate_client):
    """Process a file and upload its chunks to Weaviate."""
    try:
        # Extract content based on file type
        if file_type == 'ppt':
            chunks = extract_ppt_content(file_path)
        else:  # docx
            chunks = extract_docx_content(file_path)

        # Get the collection
        collection = weaviate_client.collections.get("Chunks")
        
        # Process each chunk
        for i, chunk in enumerate(chunks):
            # Generate embedding for the chunk
            embedding = await generate_embedding(chunk)
            if embedding is None:
                logger.error(f"Failed to generate embedding for chunk {i} in {file_path}")
                continue

            # Prepare metadata
            metadata = {
                "source_file": os.path.basename(file_path),
                "file_type": file_type,
                "chunk_index": i
            }

            # Upload to Weaviate
            try:
                collection.data.insert({
                    "content": chunk,
                    "metadata": metadata
                }, vector=embedding.tolist())
                logger.info(f"Successfully uploaded chunk {i} from {file_path}")
            except Exception as e:
                logger.error(f"Failed to upload chunk {i} from {file_path}: {str(e)}")

    except Exception as e:
        logger.error(f"Error processing file {file_path}: {str(e)}")

async def main():
    # Verify environment variables
    verify_env_variables()

    # Initialize Weaviate client
    weaviate_client = weaviate.connect_to_weaviate_cloud(
        cluster_url=WEAVIATE_URL,
        auth_credentials=AuthApiKey(WEAVIATE_API_KEY),
        skip_init_checks=True
    )

    # Get input files
    input_files = get_input_files()
    if not input_files:
        logger.warning(f"No input files found in {INPUT_DIR}")
        return

    # Process each file
    for file_path, file_type in input_files:
        logger.info(f"Processing {file_type} file: {file_path}")
        await process_and_upload_file(file_path, file_type, weaviate_client)

if __name__ == "__main__":
    asyncio.run(main())