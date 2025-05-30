import os
import asyncio
from lightrag import LightRAG
from lightrag.kg.shared_storage import initialize_pipeline_status, initialize_share_data
from lightrag.utils import setup_logger
from pptx import Presentation
from docx import Document
from typing import List, Dict, Tuple
from dotenv import load_dotenv
from lightrag.initialize import initialize_lightrag

# Load environment variables
load_dotenv()

# Set up logging
setup_logger("lightrag", level="INFO")

# Configure working directory
WORKING_DIR = "./test_data"
INPUT_DIR = os.path.join(WORKING_DIR, "PreSales")
if not os.path.exists(INPUT_DIR):
    os.makedirs(INPUT_DIR)

def verify_env_variables():
    """Verify required environment variables are set."""
    required_vars = {
        'WEAVIATE_URL': 'Weaviate vector database URL',
        'WEAVIATE_API_KEY': 'Weaviate API key',
        'NEO4J_URI': 'Neo4j database URI',
        'NEO4J_USERNAME': 'Neo4j username',
        'NEO4J_PASSWORD': 'Neo4j password'
    }
    
    missing_vars = []
    for var, description in required_vars.items():
        if not os.getenv(var):
            missing_vars.append(f"{var} ({description})")
    
    if missing_vars:
        raise ValueError(f"Missing required environment variables: {', '.join(missing_vars)}")

def extract_ppt_content(ppt_file: str) -> List[str]:
    """Extract text content from each slide."""
    prs = Presentation(ppt_file)
    slides_content = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        
        # Add slide number as header
        slide_text.append(f"[Slide {slide_num}]")
        
        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    # Add indentation to distinguish content
                    slide_text.extend("  " + line for line in text.split("\n") if line.strip())
        
        if len(slide_text) > 1:  # Only add if there's content beyond the slide number
            content = "\n".join(slide_text)
            slides_content.append(content)
    
    return slides_content

def extract_docx_content(docx_file: str) -> List[str]:
    """Extract text content from DOCX file."""
    doc = Document(docx_file)
    paragraphs_content = []
    current_chunk = []
    current_size = 0
    max_chunk_size = 2000  # Maximum characters per chunk
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
            
        # If adding this paragraph would exceed chunk size, save current chunk
        if current_size + len(text) > max_chunk_size and current_chunk:
            paragraphs_content.append("\n".join(current_chunk))
            current_chunk = []
            current_size = 0
            
        current_chunk.append(text)
        current_size += len(text)
    
    # Add any remaining content
    if current_chunk:
        paragraphs_content.append("\n".join(current_chunk))
    
    return paragraphs_content

def get_input_files() -> List[Tuple[str, str]]:
    """Get all PPT and DOCX files from input directory."""
    input_files = []
    for filename in os.listdir(INPUT_DIR):
        if filename.endswith(('.ppt', '.pptx', '.docx')):
            file_path = os.path.join(INPUT_DIR, filename)
            file_type = 'ppt' if filename.endswith(('.ppt', '.pptx')) else 'docx'
            input_files.append((file_path, file_type))
    return input_files

async def initialize_rag():
    """Initialize LightRAG with Weaviate and Neo4j."""
    # Initialize shared data and pipeline status
    initialize_share_data()
    await initialize_pipeline_status()
    
    # Set OpenAI as LLM provider
    os.environ["LLM_BINDING"] = "openai"
    os.environ["LLM_MODEL"] = "gpt-4o-mini"
    
    # Initialize LightRAG with vector and graph storage
    rag = initialize_lightrag(
        working_dir=WORKING_DIR,
        llm_provider="openai",
        llm_model_name="gpt-4o-mini",
        vector_storage="WeaviateDBVectorStorage",
        graph_storage="Neo4JStorage",
        embedding_model="sentence-transformers/all-mpnet-base-v2"
    )
    
    # Initialize storages
    await rag.initialize_storages()
    return rag

async def main():
    try:
        # Verify environment variables first
        verify_env_variables()
        
        # Initialize RAG
        print("\n🔧 Initializing LightRAG...")
        rag = await initialize_rag()
        
        # Clear document status
        await rag.doc_status.drop()
        print("✅ Document status cleared")
        print("✅ LightRAG initialized successfully")
        
        # Get all input files
        input_files = get_input_files()
        if not input_files:
            print(f"❌ No PPT or DOCX files found in {INPUT_DIR}")
            return
            
        print(f"\n📝 Found {len(input_files)} files to process")
        
        # Process each file
        for file_path, file_type in input_files:
            filename = os.path.basename(file_path)
            print(f"\n📄 Processing file: {filename}")
            
            # Extract content based on file type
            if file_type == 'ppt':
                content_list = extract_ppt_content(file_path)
                print(f"📊 Found {len(content_list)} slides with content")
            else:
                content_list = extract_docx_content(file_path)
                print(f"📄 Extracted {len(content_list)} chunks of content")
            
            # Insert content using LightRAG's pipeline methods
            print("\n⚙️ Inserting content into LightRAG...")
            for i, content in enumerate(content_list, 1):
                print(f"Processing chunk {i}/{len(content_list)}")
                
                # Create a source identifier
                source_id = f"{filename}#{'slide' if file_type == 'ppt' else 'chunk'}_{i}"
                
                # Print content for debugging
                print(f"\nContent:")
                print(content)
                print("-" * 40)
                
                await rag.apipeline_enqueue_documents(
                    input=[content],
                    file_paths=[source_id]
                )
                await rag.apipeline_process_enqueue_documents()
            
            print(f"\n✅ Successfully processed {filename}")
        
        print("\n✅ Successfully processed all files")
        
    except Exception as e:
        print(f"\n❌ An error occurred: {str(e)}")
    finally:
        if 'rag' in locals():
            await rag.finalize_storages()
            print("\n✅ Finalized storages")

if __name__ == "__main__":
    asyncio.run(main()) 